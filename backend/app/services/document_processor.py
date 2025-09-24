"""
Document processing service for extracting and preprocessing text from various file formats.
"""

import os
import re
from typing import Dict, List, Optional, Tuple
import asyncio
from pathlib import Path

import spacy
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from loguru import logger

from app.core.config import settings


class DocumentProcessor:
    """Professional document processing service."""
    
    def __init__(self):
        """Initialize the document processor."""
        self.nlp = spacy.load("en_core_web_sm")
        self.stop_words = set(stopwords.words('english'))
        self.supported_formats = ['.pdf', '.docx', '.txt', '.pptx']
        
    async def extract_text(self, file_path: str, file_type: str) -> Dict[str, any]:
        """
        Extract text from various document formats.
        
        Args:
            file_path: Path to the document file
            file_type: Type of the document (pdf, docx, txt, pptx)
            
        Returns:
            Dictionary containing raw text and metadata
        """
        try:
            if file_type.lower() == 'pdf':
                return await self._extract_pdf_text(file_path)
            elif file_type.lower() == 'docx':
                return await self._extract_docx_text(file_path)
            elif file_type.lower() == 'txt':
                return await self._extract_txt_text(file_path)
            elif file_type.lower() == 'pptx':
                return await self._extract_pptx_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    async def _extract_pdf_text(self, file_path: str) -> Dict[str, any]:
        """Extract text from PDF files."""
        def extract_sync():
            reader = PdfReader(file_path)
            text = ""
            metadata = {
                "pages": len(reader.pages),
                "title": reader.metadata.get('/Title', '') if reader.metadata else '',
                "author": reader.metadata.get('/Author', '') if reader.metadata else '',
            }
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    text += f"\n--- Page {page_num} ---\n{page_text}"
                except Exception as e:
                    logger.warning(f"Could not extract text from page {page_num}: {e}")
                    
            return {"raw_text": text.strip(), "metadata": metadata}
        
        return await asyncio.get_event_loop().run_in_executor(None, extract_sync)
    
    async def _extract_docx_text(self, file_path: str) -> Dict[str, any]:
        """Extract text from DOCX files."""
        def extract_sync():
            doc = DocxDocument(file_path)
            paragraphs = []
            headings = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Check if paragraph is a heading
                    if paragraph.style.name.startswith('Heading'):
                        headings.append({
                            "text": paragraph.text.strip(),
                            "level": int(paragraph.style.name.split()[-1])
                        })
                    paragraphs.append(paragraph.text.strip())
            
            metadata = {
                "paragraphs": len(paragraphs),
                "headings": headings
            }
            
            return {
                "raw_text": "\n\n".join(paragraphs),
                "metadata": metadata
            }
        
        return await asyncio.get_event_loop().run_in_executor(None, extract_sync)
    
    async def _extract_txt_text(self, file_path: str) -> Dict[str, any]:
        """Extract text from TXT files."""
        def extract_sync():
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            
            lines = text.split('\n')
            metadata = {
                "lines": len(lines),
                "characters": len(text)
            }
            
            return {"raw_text": text, "metadata": metadata}
        
        return await asyncio.get_event_loop().run_in_executor(None, extract_sync)
    
    async def _extract_pptx_text(self, file_path: str) -> Dict[str, any]:
        """Extract text from PPTX files."""
        def extract_sync():
            presentation = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                slides_text.append(slide_text)
            
            metadata = {
                "slides": len(presentation.slides)
            }
            
            return {
                "raw_text": "\n\n".join(slides_text),
                "metadata": metadata
            }
        
        return await asyncio.get_event_loop().run_in_executor(None, extract_sync)
    
    async def preprocess_text(self, raw_text: str) -> Dict[str, any]:
        """
        Preprocess extracted text for better analysis.
        
        Args:
            raw_text: Raw extracted text
            
        Returns:
            Dictionary containing processed text components
        """
        try:
            # Clean and normalize text
            cleaned_text = self._clean_text(raw_text)
            
            # Split into sentences
            sentences = sent_tokenize(cleaned_text)
            
            # Extract headings and sections
            sections = self._extract_sections(cleaned_text)
            
            # Process with spaCy for advanced NLP
            doc = self.nlp(cleaned_text)
            
            # Extract entities and key phrases
            entities = [(ent.text, ent.label_) for ent in doc.ents]
            key_phrases = self._extract_key_phrases(doc)
            
            # Create text chunks for processing
            chunks = self._create_chunks(sentences)
            
            return {
                "cleaned_text": cleaned_text,
                "sentences": sentences,
                "sections": sections,
                "entities": entities,
                "key_phrases": key_phrases,
                "chunks": chunks,
                "word_count": len(doc),
                "sentence_count": len(sentences)
            }
            
        except Exception as e:
            logger.error(f"Error preprocessing text: {str(e)}")
            raise
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)]', '', text)
        
        # Remove page headers/footers patterns
        text = re.sub(r'--- Page \d+ ---', '', text)
        text = re.sub(r'--- Slide \d+ ---', '', text)
        
        return text.strip()
    
    def _extract_sections(self, text: str) -> List[Dict[str, str]]:
        """Extract sections and headings from text."""
        sections = []
        lines = text.split('\n')
        
        current_section = None
        current_content = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Simple heuristic for headings (all caps, short lines, etc.)
            if (len(line) < 100 and 
                (line.isupper() or 
                 re.match(r'^[\d\.\-\s]*[A-Z]', line) or
                 line.endswith(':'))):
                
                # Save previous section
                if current_section:
                    sections.append({
                        "heading": current_section,
                        "content": " ".join(current_content)
                    })
                
                # Start new section
                current_section = line
                current_content = []
            else:
                if current_section:
                    current_content.append(line)
        
        # Add final section
        if current_section and current_content:
            sections.append({
                "heading": current_section,
                "content": " ".join(current_content)
            })
        
        return sections
    
    def _extract_key_phrases(self, doc) -> List[str]:
        """Extract key phrases using spaCy."""
        key_phrases = []
        
        # Extract noun phrases
        for chunk in doc.noun_chunks:
            if len(chunk.text.split()) > 1 and chunk.text.lower() not in self.stop_words:
                key_phrases.append(chunk.text)
        
        # Extract named entities
        for ent in doc.ents:
            if ent.label_ in ['PERSON', 'ORG', 'GPE', 'PRODUCT', 'EVENT']:
                key_phrases.append(ent.text)
        
        return list(set(key_phrases))
    
    def _create_chunks(self, sentences: List[str], chunk_size: int = 5) -> List[Dict[str, any]]:
        """Create overlapping text chunks for processing."""
        chunks = []
        
        for i in range(0, len(sentences), chunk_size - 1):
            chunk_sentences = sentences[i:i + chunk_size]
            chunk_text = " ".join(chunk_sentences)
            
            chunks.append({
                "id": i // (chunk_size - 1),
                "text": chunk_text,
                "start_sentence": i,
                "end_sentence": min(i + chunk_size - 1, len(sentences) - 1),
                "sentence_count": len(chunk_sentences)
            })
            
            if i + chunk_size >= len(sentences):
                break
        
        return chunks


# Global instance
document_processor = DocumentProcessor()